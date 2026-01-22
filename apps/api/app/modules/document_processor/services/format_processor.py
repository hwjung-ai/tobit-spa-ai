"""Format processor for multi-format document extraction"""

import logging
from typing import Any, Dict

try:
    import pdf2image
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
    pdf2image = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
    import pandas as pd
except ImportError:
    openpyxl = None
    pd = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    Image = None
    pytesseract = None

logger = logging.getLogger(__name__)


class DocumentProcessingError(Exception):
    """Exception raised during document processing"""
    pass


class DocumentProcessor:
    """Main processor for handling various document formats"""

    SUPPORTED_FORMATS = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'png': 'image/png',
    }

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def process(self, file_path: str, format: str) -> Dict[str, Any]:
        """
        Main entry point for document processing

        Args:
            file_path: Path to the file
            format: File format (pdf, docx, xlsx, pptx, image)

        Returns:
            Dictionary with extracted content and metadata
        """

        format = format.lower()

        if format == "pdf":
            return self._process_pdf(file_path)
        elif format == "docx":
            return self._process_docx(file_path)
        elif format == "xlsx":
            return self._process_xlsx(file_path)
        elif format == "pptx":
            return self._process_pptx(file_path)
        elif format in ["jpg", "jpeg", "png"]:
            return self._process_image(file_path)
        else:
            raise DocumentProcessingError(f"Unsupported format: {format}")

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PDF"""

        if PdfReader is None:
            raise DocumentProcessingError("pypdf library not installed")

        try:
            pdf_document = PdfReader(file_path)
            pages = []

            for page_num, page in enumerate(pdf_document.pages):
                # Try direct text extraction
                text = page.extract_text()

                # If extraction failed or returned very little text, we could try OCR
                # but for now, just use what we got
                if not text:
                    text = ""

                pages.append({
                    "page_num": page_num,
                    "text": text,
                    "tables": [],
                    "images": []
                })

            return {
                "pages": pages,
                "metadata": {
                    "page_count": len(pages),
                    "format": "pdf",
                    "extraction_method": "pdf_text"
                }
            }

        except Exception as e:
            raise DocumentProcessingError(f"PDF processing failed: {str(e)}")

    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract content from DOCX"""

        if DocxDocument is None:
            raise DocumentProcessingError("python-docx library not installed")

        try:
            doc = DocxDocument(file_path)
            content = {
                "pages": [{
                    "page_num": 0,
                    "text": "",
                    "tables": [],
                    "images": []
                }],
                "metadata": {
                    "format": "docx",
                    "extraction_method": "docx_native"
                }
            }

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    content["pages"][0]["text"] += paragraph.text + "\n"

            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text for cell in row.cells])

                content["pages"][0]["tables"].append({
                    "data": table_data,
                    "rows": len(table.rows),
                    "cols": len(table.columns) if table.columns else 0
                })

            return content

        except Exception as e:
            raise DocumentProcessingError(f"DOCX processing failed: {str(e)}")

    def _process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Extract content from Excel"""

        if pd is None:
            raise DocumentProcessingError("pandas library not installed")

        try:
            excel_file = pd.ExcelFile(file_path)
            sheets = []

            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheets.append({
                    "sheet_name": sheet_name,
                    "columns": df.columns.tolist(),
                    "rows": len(df),
                    "data_preview": df.head(10).to_dict('records') if len(df) > 0 else []
                })

            content = {
                "pages": [{
                    "page_num": 0,
                    "text": f"Excel file with {len(sheets)} sheets",
                    "tables": sheets,
                    "images": []
                }],
                "metadata": {
                    "format": "xlsx",
                    "extraction_method": "xlsx_native",
                    "sheet_count": len(sheets)
                }
            }

            return content

        except Exception as e:
            raise DocumentProcessingError(f"XLSX processing failed: {str(e)}")

    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PowerPoint"""

        if Presentation is None:
            raise DocumentProcessingError("python-pptx library not installed")

        try:
            prs = Presentation(file_path)
            slides = []

            for slide_num, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_num": slide_num,
                    "text": "",
                    "shapes": []
                }

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_data["text"] += shape.text + "\n"
                        slide_data["shapes"].append({
                            "type": "text",
                            "content": shape.text
                        })
                    elif hasattr(shape, "table"):
                        # Extract table
                        table_data = []
                        for row in shape.table.rows:
                            table_data.append([cell.text for cell in row.cells])

                        slide_data["shapes"].append({
                            "type": "table",
                            "data": table_data
                        })

                slides.append(slide_data)

            return {
                "pages": slides,
                "metadata": {
                    "format": "pptx",
                    "extraction_method": "pptx_native",
                    "slide_count": len(slides)
                }
            }

        except Exception as e:
            raise DocumentProcessingError(f"PPTX processing failed: {str(e)}")

    def _process_image(self, file_path: str) -> Dict[str, Any]:
        """Extract text from image using OCR"""

        if Image is None or pytesseract is None:
            raise DocumentProcessingError("PIL or pytesseract library not installed")

        try:
            image = Image.open(file_path)

            # Try pytesseract OCR
            text = pytesseract.image_to_string(image)

            return {
                "pages": [{
                    "page_num": 0,
                    "text": text,
                    "tables": [],
                    "images": []
                }],
                "metadata": {
                    "format": "image",
                    "extraction_method": "ocr_pytesseract",
                    "width": image.width,
                    "height": image.height
                }
            }

        except Exception as e:
            raise DocumentProcessingError(f"Image processing failed: {str(e)}")
